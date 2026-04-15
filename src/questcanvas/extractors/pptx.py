"""PPTX extraction via python-pptx."""

from __future__ import annotations

from io import BytesIO

from ..canvas.types import CanvasFile
from ..errors import ExtractionError, MissingDependencyError
from .base import ExtractionUnit, ExtractionResult, build_result


class PptxExtractor:
    """Extractor for PowerPoint files with slide markers."""

    def supports(self, canvas_file: CanvasFile) -> bool:
        content_type = (canvas_file.content_type or "").lower()
        return (
            content_type
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            or canvas_file.effective_name.lower().endswith(".pptx")
        )

    def extract(
        self,
        canvas_file: CanvasFile,
        content: bytes,
        *,
        max_units: int,
    ) -> ExtractionResult:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise MissingDependencyError(
                "python-pptx is required for slide extraction. Install questcanvas dependencies first."
            ) from exc

        try:
            presentation = Presentation(BytesIO(content))
        except Exception as exc:
            raise ExtractionError(f"Could not open PPTX file {canvas_file.effective_name}.") from exc

        units: list[ExtractionUnit] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            text_parts: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if isinstance(text, str):
                    cleaned = text.strip()
                    if cleaned:
                        text_parts.append(cleaned)
            if not text_parts:
                continue
            units.append(
                ExtractionUnit(
                    unit_type="slide",
                    unit_number=slide_number,
                    text=f"[Slide {slide_number}]\n" + "\n".join(text_parts),
                )
            )

        return build_result(canvas_file, unit_type="slide", units=units, max_units=max_units)
